import re
from pathlib import Path


def render_document_markdown(document: dict) -> str:
    title = str(document.get("title") or "Document").strip() or "Document"
    topic = str(document.get("topic") or "").strip()
    summary = str(document.get("summary") or "").strip()
    sections = document.get("sections") if isinstance(document.get("sections"), list) else []
    extension_reading = document.get("extension_reading") if isinstance(document.get("extension_reading"), list) else []

    lines = [f"# {title}", ""]
    if topic:
        lines.extend([f"Topic: {topic}", ""])
    if summary:
        lines.extend([summary, ""])

    for section in sections:
        if not isinstance(section, dict):
            continue
        heading = str(section.get("heading") or "").strip()
        body = str(section.get("body") or "").strip()
        if heading:
            lines.extend([f"## {heading}", ""])
        lines.extend([body or "N/A", ""])
        structured_blocks = [
            ("Key Points", section.get("key_points")),
            ("Examples", section.get("examples")),
            ("Common Pitfalls", section.get("pitfalls")),
            ("Self Check", section.get("checklist")),
            ("Evidence", section.get("evidence")),
        ]
        for block_title, raw_items in structured_blocks:
            items = raw_items if isinstance(raw_items, list) else []
            normalized_items = [str(item or "").strip() for item in items if str(item or "").strip()]
            if not normalized_items:
                continue
            lines.extend([f"### {block_title}", ""])
            for item in normalized_items:
                lines.append(f"- {item}")
            lines.append("")

    if extension_reading:
        lines.extend(["## Extension Reading", ""])
        for item in extension_reading:
            if not isinstance(item, dict):
                continue
            item_title = str(item.get("title") or "").strip() or "Untitled"
            item_reason = str(item.get("reason") or "").strip()
            lines.append(f"- {item_title}")
            if item_reason:
                lines.append(f"  - Reason: {item_reason}")
        lines.append("")

    return "\n".join(lines).strip() + "\n"


def render_quiz_markdown(quiz: dict) -> str:
    title = str(quiz.get("title") or "Quiz").strip() or "Quiz"
    topic = str(quiz.get("topic") or "").strip()
    questions = quiz.get("questions") if isinstance(quiz.get("questions"), list) else []

    lines = [f"# {title}", ""]
    if topic:
        lines.extend([f"Topic: {topic}", ""])

    for index, question in enumerate(questions, start=1):
        if not isinstance(question, dict):
            continue
        q_type = str(question.get("type") or "unknown").strip() or "unknown"
        difficulty = str(question.get("difficulty") or "").strip()
        stem = str(question.get("stem") or "").strip()
        answer = question.get("answer")
        explanation = str(question.get("explanation") or "").strip()
        knowledge_points = question.get("knowledge_points") if isinstance(question.get("knowledge_points"), list) else []

        heading = f"## Q{index}. {q_type}"
        if difficulty:
            heading += f" ({difficulty})"
        lines.extend([heading, "", stem or "No stem provided.", ""])

        if q_type == "single_choice" and isinstance(question.get("options"), list):
            for opt_index, option in enumerate(question["options"]):
                label = chr(ord("A") + opt_index)
                lines.append(f"- {label}. {option}")
            lines.append("")

        answer_text = "True" if answer is True else "False" if answer is False else str(answer)
        lines.extend([f"Answer: {answer_text}", ""])
        lines.extend([f"Explanation: {explanation or 'N/A'}", ""])

        if knowledge_points:
            lines.extend([f"Knowledge Points: {', '.join(map(str, knowledge_points))}", ""])

    return "\n".join(lines).strip() + "\n"


def render_coding_practice_markdown(practice: dict) -> str:
    title = str(practice.get("title") or "Coding Practice").strip() or "Coding Practice"
    topic = str(practice.get("topic") or "").strip()
    summary = str(practice.get("summary") or "").strip()
    language = str(practice.get("language") or "").strip() or "text"
    learning_objectives = (
        practice.get("learning_objectives")
        if isinstance(practice.get("learning_objectives"), list)
        else []
    )
    steps = practice.get("steps") if isinstance(practice.get("steps"), list) else []
    code_files = practice.get("code_files") if isinstance(practice.get("code_files"), list) else []
    run_guide = practice.get("run_guide") if isinstance(practice.get("run_guide"), dict) else {}

    lines = [f"# {title}", ""]
    if topic:
        lines.extend([f"Topic: {topic}", ""])
    if summary:
        lines.extend([summary, ""])

    if learning_objectives:
        lines.extend(["## Learning Objectives", ""])
        for item in learning_objectives:
            text = str(item or "").strip()
            if text:
                lines.append(f"- {text}")
        lines.append("")

    if steps:
        lines.extend(["## Practice Steps", ""])
        for index, step in enumerate(steps, start=1):
            if not isinstance(step, dict):
                continue
            step_title = str(step.get("title") or f"Step {index}").strip() or f"Step {index}"
            instruction = str(step.get("instruction") or "").strip()
            lines.extend([f"### {index}. {step_title}", ""])
            lines.extend([instruction or "N/A", ""])

    if run_guide:
        lines.extend(["## Run Guide", ""])
        entry_file = str(run_guide.get("entry_file") or "").strip()
        command = str(run_guide.get("command") or "").strip()
        expected_output = str(run_guide.get("expected_output") or "").strip()
        if entry_file:
            lines.extend([f"Entry File: `{entry_file}`", ""])
        if command:
            lines.extend([f"Command: `{command}`", ""])
        if expected_output:
            lines.extend(["Expected Output:", "", f"```text\n{expected_output}\n```", ""])

    if code_files:
        lines.extend(["## Code Files", ""])
        for item in code_files:
            if not isinstance(item, dict):
                continue
            file_path = str(item.get("path") or "").strip() or "code/main.py"
            purpose = str(item.get("purpose") or "").strip()
            content = str(item.get("content") or "").rstrip()
            heading = f"### `{file_path}`"
            if purpose:
                heading += f" ({purpose})"
            lines.extend([heading, ""])
            lines.extend([f"```{language}", content or "# empty", "```", ""])

    return "\n".join(lines).strip() + "\n"


def render_ppt_markdown(ppt: dict) -> str:
    title = str(ppt.get("title") or "PPT").strip() or "PPT"
    topic = str(ppt.get("topic") or "").strip()
    summary = str(ppt.get("summary") or "").strip()
    theme = str(ppt.get("theme") or "").strip()
    slide_style = str(ppt.get("slide_style") or "").strip()
    slides = ppt.get("slides") if isinstance(ppt.get("slides"), list) else []

    lines = [f"# {title}", ""]
    if topic:
        lines.extend([f"Topic: {topic}", ""])
    if theme:
        lines.extend([f"Theme: {theme}", ""])
    if slide_style:
        lines.extend([f"Style: {slide_style}", ""])
    if summary:
        lines.extend([summary, ""])

    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        slide_title = str(slide.get("title") or f"Slide {index}").strip() or f"Slide {index}"
        body = str(slide.get("body") or "").strip()
        bullets = slide.get("bullets") if isinstance(slide.get("bullets"), list) else []
        speaker_notes = str(slide.get("speaker_notes") or "").strip()
        visual_hint = str(slide.get("visual_hint") or "").strip()

        lines.extend([f"## Slide {index}: {slide_title}", ""])
        if body:
            lines.extend([body, ""])
        for bullet in bullets:
            bullet_text = str(bullet or "").strip()
            if bullet_text:
                lines.append(f"- {bullet_text}")
        if bullets:
            lines.append("")
        if visual_hint:
            lines.extend([f"Visual Hint: {visual_hint}", ""])
        if speaker_notes:
            lines.extend(["Speaker Notes:", "", speaker_notes, ""])

    return "\n".join(lines).strip() + "\n"


def _coerce_rgb(color_module, value, fallback):
    try:
        text = str(value or "").strip().lstrip("#")
        if len(text) == 6:
            return color_module(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except Exception:
        pass
    return fallback


def _pick_ppt_palette(theme: str, color_module):
    palettes = {
        "academic-clean": {
            "background": color_module(255, 255, 255),
            "title": color_module(22, 41, 70),
            "body": color_module(45, 52, 64),
            "accent": color_module(85, 112, 154),
            "accent_soft": color_module(226, 235, 247),
            "muted": color_module(116, 127, 141),
            "divider": color_module(218, 225, 234),
            "cover_background": color_module(243, 247, 252),
            "cover_panel": color_module(35, 62, 102),
            "cover_text": color_module(255, 255, 255),
            "callout_background": color_module(248, 250, 253),
        },
        "academic-rich": {
            "background": color_module(249, 247, 242),
            "title": color_module(38, 30, 20),
            "body": color_module(66, 55, 40),
            "accent": color_module(174, 93, 49),
            "accent_soft": color_module(246, 230, 219),
            "muted": color_module(127, 111, 96),
            "divider": color_module(229, 214, 198),
            "cover_background": color_module(252, 244, 236),
            "cover_panel": color_module(133, 67, 33),
            "cover_text": color_module(255, 247, 240),
            "callout_background": color_module(255, 251, 246),
        },
        "contrast-blue": {
            "background": color_module(244, 248, 255),
            "title": color_module(10, 36, 73),
            "body": color_module(32, 56, 88),
            "accent": color_module(0, 109, 219),
            "accent_soft": color_module(220, 236, 255),
            "muted": color_module(96, 119, 149),
            "divider": color_module(207, 223, 242),
            "cover_background": color_module(233, 243, 255),
            "cover_panel": color_module(0, 81, 165),
            "cover_text": color_module(255, 255, 255),
            "callout_background": color_module(240, 247, 255),
        },
    }
    return dict(palettes.get(theme, palettes["academic-clean"]))


def _pt(base_size: int, *, extra: int = 10) -> int:
    return base_size + extra


def _normalize_bullet_items(bullets):
    items = []
    if not isinstance(bullets, list):
        return items
    for bullet in bullets:
        raw_text = str(bullet or "").strip()
        if not raw_text:
            continue
        parts = [part.strip(" -•\t") for part in re.split(r"[\r\n]+", raw_text) if part.strip()]
        items.extend(part for part in parts if part)
    return items


def _normalize_body_paragraphs(value):
    if isinstance(value, list):
        candidates = value
    else:
        text = str(value or "").strip()
        if not text:
            return []
        candidates = re.split(r"[\r\n]+", text)
    paragraphs = []
    for item in candidates:
        text = " ".join(str(item or "").split())
        if text:
            paragraphs.append(text)
    return paragraphs[:3]


def _strip_step_prefix(text: str) -> str:
    return re.sub(r"^\s*(?:第\s*)?\d+\s*(?:步)?[\.\)、:：-]?\s*", "", text).strip()


def _extract_label_value(text: str):
    for separator in ("：", ":"):
        if separator not in text:
            continue
        left, right = text.split(separator, 1)
        left = left.strip(" -•\t")
        right = right.strip()
        if left and right and len(left) <= 18:
            return left, right
    return None


def _split_detail_fragments(text: str):
    fragments = [fragment.strip() for fragment in re.split(r"[；;]\s*", text) if fragment.strip()]
    if len(fragments) >= 2:
        return fragments[:3]
    return [text.strip()] if text.strip() else []


def _expand_bullet_text(text: str) -> dict:
    cleaned = _strip_step_prefix(str(text or "").strip())
    if not cleaned:
        return {"headline": "", "details": []}

    label_value = _extract_label_value(cleaned)
    if label_value:
        headline, detail = label_value
        return {
            "headline": headline,
            "details": _split_detail_fragments(detail),
        }

    segments = [segment.strip() for segment in re.split(r"[；;]\s*", cleaned) if segment.strip()]
    if len(segments) >= 2:
        return {"headline": segments[0], "details": segments[1:3]}

    for separator in ("。", "，", ","):
        if separator not in cleaned:
            continue
        headline, detail = cleaned.split(separator, 1)
        headline = headline.strip()
        detail = detail.strip()
        if headline and detail and 4 <= len(headline) <= 20:
            return {"headline": headline, "details": [detail]}

    return {"headline": cleaned, "details": []}


def _extract_table_rows(items):
    rows = []
    remainder = []
    for item in items:
        label_value = _extract_label_value(_strip_step_prefix(item))
        if label_value:
            rows.append(label_value)
        else:
            remainder.append(item)
    if len(rows) >= 3 and len(rows) >= len(remainder) + 2:
        return rows[:5], remainder[:1]
    return [], items


def _build_process_steps(items):
    steps = []
    for index, item in enumerate(items, start=1):
        cleaned = str(item or "").strip()
        match = re.match(r"^\s*(?:第\s*)?(\d+)\s*(?:步)?[\.\)、:：-]?\s*", cleaned)
        step_number = int(match.group(1)) if match else index
        expanded = _expand_bullet_text(cleaned)
        headline = expanded["headline"] or f"步骤 {step_number}"
        steps.append(
            {
                "number": step_number,
                "headline": headline,
                "details": expanded["details"][:2],
            }
        )
    return steps[:4]


def _infer_slide_role(slide_payload: dict, index: int, total_slides: int) -> str:
    title = str(slide_payload.get("title") or "").strip().lower()
    if index == 1 or "封面" in title:
        return "cover"
    if "总结" in title or "回顾" in title:
        return "summary"
    return "content"


def render_pptx_file(ppt: dict, output_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required to export ppt resources as .pptx files"
        ) from exc

    title = str(ppt.get("title") or "PPT").strip() or "PPT"
    topic = str(ppt.get("topic") or "").strip()
    summary = str(ppt.get("summary") or "").strip()
    theme = str(ppt.get("theme") or "").strip().lower()
    slides = ppt.get("slides") if isinstance(ppt.get("slides"), list) else []
    palette = _pick_ppt_palette(theme or "academic-clean", RGBColor)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = title
    if topic:
        presentation.core_properties.subject = topic
    if summary:
        presentation.core_properties.comments = summary

    total_slides = max(1, len(slides))
    blank_layout = presentation.slide_layouts[6]

    def _style_paragraph(paragraph, size_pt: int, color, *, bold: bool = False, align=None):
        if align is not None:
            paragraph.alignment = align
        paragraph.space_after = Pt(4)
        for run in paragraph.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color

    def _set_shape_fill(shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    def _clip_text(text: str, max_chars: int) -> str:
        compact = " ".join(str(text or "").split())
        if len(compact) <= max_chars:
            return compact
        return compact[: max(0, max_chars - 1)].rstrip() + "…"

    def _fit_font_size(items, role_name: str) -> int:
        line_count = 0
        for item in items[:6]:
            expanded = _expand_bullet_text(item)
            if expanded["headline"]:
                line_count += 1
            line_count += min(2, len(expanded["details"]))
        if role_name == "cover":
            return _pt(14)
        if line_count >= 10 or len(items) >= 6:
            return _pt(11)
        if line_count >= 7 or len(items) >= 4:
            return _pt(12)
        return _pt(14)

    def _render_bullet_list(slide, items, *, left: float, top: float, width: float, height: float, role_name: str, color, detail_color):
        bullet_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = bullet_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        body_size = _fit_font_size(items, role_name)
        detail_size = max(_pt(8), body_size - 4)
        first_paragraph = True
        max_items = 5 if role_name != "cover" else 4
        for bullet in items[:max_items]:
            expanded = _expand_bullet_text(bullet)
            headline = _clip_text(expanded["headline"], 34 if width < 5.5 else 60)
            if not headline:
                continue
            paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
            first_paragraph = False
            paragraph.text = f"• {headline}"
            paragraph.level = 0
            _style_paragraph(paragraph, body_size, color, bold=True, align=PP_ALIGN.LEFT)
            for detail in expanded["details"][:2]:
                detail_paragraph = text_frame.add_paragraph()
                detail_paragraph.text = _clip_text(detail, 42 if width < 5.5 else 84)
                detail_paragraph.level = 1
                _style_paragraph(detail_paragraph, detail_size, detail_color, align=PP_ALIGN.LEFT)
        return bullet_box

    for index, slide_payload in enumerate(slides, start=1):
        if not isinstance(slide_payload, dict):
            continue

        slide = presentation.slides.add_slide(blank_layout)
        role = _infer_slide_role(slide_payload, index, total_slides)
        accent_color = _coerce_rgb(RGBColor, slide_payload.get("accent_color"), palette["accent"])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = palette["cover_background"] if role == "cover" else palette["background"]

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.5), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = str(slide_payload.get("title") or f"Slide {index}").strip() or f"Slide {index}"
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_run = title_paragraph.runs[0]
        title_run.font.name = "Aptos Display" if role == "cover" else "Aptos"
        title_run.font.size = Pt(_pt(28 if role == "cover" else 24))
        title_run.font.bold = True
        title_run.font.color.rgb = palette["cover_text"] if role == "cover" else palette["title"]

        divider = slide.shapes.add_shape(1, Inches(0.72), Inches(1.22), Inches(11.55), Inches(0.02))
        divider.fill.solid()
        divider.fill.fore_color.rgb = accent_color if role == "cover" else palette["divider"]
        divider.line.fill.background()

        bullets = _normalize_bullet_items(slide_payload.get("bullets"))
        body_paragraphs = _normalize_body_paragraphs(slide_payload.get("body"))
        left_bullets = bullets
        right_bullets = []
        table_rows = []
        table_remainder = []
        process_steps = []

        if role == "cover":
            cover_panel = slide.shapes.add_shape(1, Inches(0.95), Inches(1.62), Inches(11.4), Inches(3.9))
            cover_panel.fill.solid()
            cover_panel.fill.fore_color.rgb = palette["cover_panel"]
            cover_panel.line.color.rgb = accent_color

            subtitle_box = slide.shapes.add_textbox(Inches(1.35), Inches(2.02), Inches(10.55), Inches(1.2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            subtitle_paragraph.text = _clip_text(summary or topic or "课程导入", 104)
            if subtitle_paragraph.runs:
                subtitle_run = subtitle_paragraph.runs[0]
                subtitle_run.font.name = "Aptos"
                subtitle_run.font.size = Pt(_pt(15))
                subtitle_run.font.color.rgb = palette["cover_text"]
            if bullets:
                _render_bullet_list(
                    slide,
                    bullets,
                    left=1.7,
                    top=3.48,
                    width=9.85,
                    height=1.18,
                    role_name="cover",
                    color=palette["cover_text"],
                    detail_color=palette["cover_text"],
                )
        else:
            content_panel = slide.shapes.add_shape(1, Inches(0.92), Inches(1.5), Inches(11.5), Inches(5.05))
            content_panel.fill.solid()
            content_panel.fill.fore_color.rgb = palette["callout_background"]
            content_panel.line.color.rgb = palette["divider"]

            if body_paragraphs:
                body_box = slide.shapes.add_textbox(Inches(1.18), Inches(1.78), Inches(10.9), Inches(1.0))
                body_frame = body_box.text_frame
                body_frame.clear()
                body_frame.word_wrap = True
                body_frame.margin_left = 0
                body_frame.margin_right = 0
                body_paragraph = body_frame.paragraphs[0]
                body_paragraph.text = _clip_text(body_paragraphs[0], 120)
                _style_paragraph(body_paragraph, _pt(10), palette["body"], align=PP_ALIGN.LEFT)
            _render_bullet_list(
                slide,
                bullets,
                left=1.18,
                top=3.0 if body_paragraphs else 1.82,
                width=10.9,
                height=3.35 if body_paragraphs else 4.5,
                role_name=role,
                color=palette["body"],
                detail_color=palette["muted"],
            )

        footer_line = slide.shapes.add_shape(1, Inches(0.92), Inches(6.86), Inches(11.5), Inches(0.02))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = palette["divider"] if role != "cover" else accent_color
        footer_line.line.fill.background()

        footer_box = slide.shapes.add_textbox(Inches(0.95), Inches(6.94), Inches(11.35), Inches(0.18))
        footer_frame = footer_box.text_frame
        footer_frame.clear()
        footer_paragraph = footer_frame.paragraphs[0]
        footer_paragraph.alignment = PP_ALIGN.RIGHT
        footer_paragraph.text = f"{topic or title}  |  {index}/{total_slides}"
        if footer_paragraph.runs:
            footer_run = footer_paragraph.runs[0]
            footer_run.font.name = "Aptos"
            footer_run.font.size = Pt(_pt(8))
            footer_run.font.color.rgb = palette["cover_text"] if role == "cover" else palette["muted"]

        speaker_notes = str(slide_payload.get("speaker_notes") or "").strip()
        if speaker_notes:
            try:
                notes_slide = slide.notes_slide
                notes_frame = getattr(notes_slide, "notes_text_frame", None)
                if notes_frame is not None:
                    notes_frame.text = speaker_notes
                    for paragraph in notes_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(_pt(12))
            except Exception:
                pass

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
